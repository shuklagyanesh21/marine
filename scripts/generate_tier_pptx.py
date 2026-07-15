#!/usr/bin/env python3
"""Generate a comprehensive PowerPoint explaining the tier-wise genome download
process and the SmORFinder hybrid run (appended as the final three slides).
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color palette ───────────────────────────────────────────────────────────
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
OCEAN_BLUE = RGBColor(0x14, 0x6E, 0xB4)
TEAL = RGBColor(0x0E, 0x91, 0x8C)
DARK_TEAL = RGBColor(0x09, 0x6B, 0x67)
LIGHT_BLUE = RGBColor(0xD6, 0xEC, 0xF5)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x1B, 0x7A, 0x2E)
RED = RGBColor(0xC0, 0x39, 0x2B)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_box(slide, text, left, top, width, height, font_size=32, bold=True, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return tf


def add_body_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet(tf, text, level=0, font_size=16, bold=False, color=DARK_GRAY, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = space_before
    return p


def add_table(slide, rows, col_widths, left, top, row_height=Inches(0.4)):
    n_rows = len(rows)
    n_cols = len(rows[0])
    width = sum(col_widths)
    height = row_height * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = OCEAN_BLUE
            else:
                p.font.color.rgb = DARK_GRAY
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BLUE

    return table_shape


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    add_title_box(slide, "Marine Prokaryote Genome Collection",
                  Inches(1), Inches(1.5), Inches(11), Inches(1.2), font_size=40, color=WHITE)
    add_title_box(slide, "Tier-wise Download Pipeline",
                  Inches(1), Inches(2.7), Inches(11), Inches(0.8), font_size=28, bold=False, color=LIGHT_BLUE)

    tf = add_body_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(2.5))
    add_bullet(tf, "~120,000 marine prokaryote genomes across 3 tiers", font_size=18, color=WHITE)
    add_bullet(tf, "Source: GTDB R220 + external curated catalogs", font_size=18, color=WHITE)
    add_bullet(tf, "Goal: Discover novel antimicrobial peptides from marine microbes", font_size=18, color=WHITE)
    add_bullet(tf, "", font_size=10, color=WHITE)
    add_bullet(tf, "Marine Peptides Project | June 2026", font_size=14, color=MID_GRAY)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 2: Project Overview
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Project Overview", Inches(0.7), Inches(0.3), Inches(11), Inches(0.8),
                  font_size=30, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "The Challenge"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, '"Marine" is not a taxonomy — no phylum/class/kingdom called marine', color=DARK_GRAY)
    add_bullet(tf, "Marine origin is buried in free-text metadata fields (inconsistent, often missing)", color=DARK_GRAY)
    add_bullet(tf, "Solution: a 3-tier progressive recall strategy", color=DARK_GRAY)

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Our Approach", font_size=20, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "Tier 1: Offline keyword filtering of GTDB metadata (high precision, free)", color=DARK_GRAY)
    add_bullet(tf, "Tier 2: NCBI BioSample ENVO attribute queries (recovers hidden marine MAGs)", color=DARK_GRAY)
    add_bullet(tf, "Tier 3: External curated marine catalogs (maximum recall)", color=DARK_GRAY)

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Key Principles", font_size=20, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "Config-driven (config/config.yaml) — all keywords, thresholds, paths in one place", color=DARK_GRAY)
    add_bullet(tf, "Reproducible — re-run 5 scripts to regenerate everything from metadata", color=DARK_GRAY)
    add_bullet(tf, "Resumable — every download step can be interrupted and resumed safely", color=DARK_GRAY)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 3: Pipeline Architecture
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Pipeline Architecture", Inches(0.7), Inches(0.3), Inches(11), Inches(0.8),
                  font_size=30, color=NAVY)

    # Flow diagram using shapes
    shapes_data = [
        ("GTDB R220\nAll prokaryotes\n~400k genomes", Inches(5), Inches(1.3), Inches(3.3), Inches(0.9), NAVY),
        ("Tier 1\nGTDB keyword filter\n→ 37,069 genomes", Inches(1.5), Inches(2.8), Inches(3.3), Inches(0.9), OCEAN_BLUE),
        ("Tier 2\nBioSample ENVO\n→ 5,030 genomes", Inches(5), Inches(2.8), Inches(3.3), Inches(0.9), TEAL),
        ("Tier 3\nExternal catalogs\n→ 84,934 genomes", Inches(8.5), Inches(2.8), Inches(3.3), Inches(0.9), DARK_TEAL),
        ("Combined: ~120,396 marine genomes", Inches(3.5), Inches(4.5), Inches(6.3), Inches(0.7), GREEN),
    ]

    for text, left, top, width, height, color in shapes_data:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate(text.split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(12) if i > 0 else Pt(14)
            p.font.bold = (i == 0)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    # Summary table
    rows_data = [
        ["Tier", "Method", "Genomes", "Precision", "Recall"],
        ["1", "GTDB keyword filter (offline)", "37,069", "High", "Moderate"],
        ["2", "NCBI BioSample ENVO queries", "5,030", "High", "High for MAGs"],
        ["3", "External curated catalogs", "84,934", "Moderate", "Maximum"],
    ]
    add_table(slide, rows_data,
              [Inches(0.6), Inches(4.0), Inches(1.4), Inches(1.4), Inches(2.0)],
              Inches(1.5), Inches(5.6))

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 4: Tier 1 — GTDB Keyword Filter
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Tier 1 — GTDB Metadata Keyword Filter",
                  Inches(0.7), Inches(0.3), Inches(11), Inches(0.8), font_size=30, color=NAVY)

    # Left column
    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(6), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "How it works"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "Parse GTDB R220 bac120 + ar53 metadata TSVs (~400k genomes)")
    add_bullet(tf, "Scan fields: isolation_source, country, organism_name, lat_lon")
    add_bullet(tf, "Match marine include keywords (whole-word, case-insensitive)")
    add_bullet(tf, "Reject if exclusion keywords hit (freshwater, soil, clinical…)")
    add_bullet(tf, "Tag host-associated genomes (sponge, coral, fish…)")

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Result: 37,069 genomes (18,371 species representatives)", bold=True, color=GREEN)

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Script", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "scripts/00_select_marine_candidates.py")
    add_bullet(tf, "Logic: src/marine_peptides/download/marine_filter.py")

    # Right column - keyword categories
    tf2 = add_body_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf2.paragraphs[0]
    p.text = "Keyword Categories"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf2, "Core: marine, seawater, ocean, oceanic", font_size=14)
    add_bullet(tf2, "Physical: deep sea, abyssal, hadal, pelagic", font_size=14)
    add_bullet(tf2, "Features: hydrothermal vent, cold seep", font_size=14)
    add_bullet(tf2, "Coastal: estuary, intertidal, mangrove, reef", font_size=14)
    add_bullet(tf2, "Geographic: Pacific, Atlantic, Mediterranean…", font_size=14)

    add_bullet(tf2, "", font_size=8)
    add_bullet(tf2, "Exclusion Terms", font_size=18, bold=True, color=RED)
    add_bullet(tf2, "freshwater, lake, river, soil, terrestrial", font_size=14)
    add_bullet(tf2, "wastewater, sewage, sludge", font_size=14)
    add_bullet(tf2, "gut, clinical, hospital, blood", font_size=14)
    add_bullet(tf2, "Marine Drive, Marine Corps (false positives)", font_size=14)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 5: Tier 2 — BioSample ENVO
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Tier 2 — NCBI BioSample Enrichment",
                  Inches(0.7), Inches(0.3), Inches(11), Inches(0.8), font_size=30, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(6), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "Why Tier 2?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "~80% of GTDB isolation_source is EMPTY")
    add_bullet(tf, "Many marine MAGs/SAGs have env data only in NCBI BioSample records")
    add_bullet(tf, "These are invisible to GTDB-only filtering")

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Method", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "Filter to environmental genomes (MAG/SAG/env) with a BioSample ID")
    add_bullet(tf, "Exclude those already flagged by Tier 1 → 246,025 candidates")
    add_bullet(tf, "Query NCBI E-utilities (esearch → efetch by UID) in batches of 200")
    add_bullet(tf, "Rate-limited: 3 req/s without API key, 10 req/s with NCBI_API_KEY")
    add_bullet(tf, "Search BioSample fields: env_broad_scale, env_local_scale,")
    add_bullet(tf, "  env_medium, isolation_source, geo_loc_name", level=1, font_size=14)
    add_bullet(tf, "Apply same marine keyword list → match = marine")

    tf2 = add_body_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf2.paragraphs[0]
    p.text = "Results"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf2, "5,030 additional marine genomes recovered", bold=True, color=GREEN)
    add_bullet(tf2, "2,691 are GTDB species representatives")
    add_bullet(tf2, "")
    add_bullet(tf2, "Evidence strength:", bold=True)
    add_bullet(tf2, "  83% strong signal (env_broad/medium/isolation)")
    add_bullet(tf2, "  2% weaker (geo_loc_name only)")
    add_bullet(tf2, "  Every match recorded in marine_evidence column")

    add_bullet(tf2, "", font_size=8)
    add_bullet(tf2, "Scripts", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf2, "scripts/01_enrich_marine_biosample.py", font_size=14)
    add_bullet(tf2, "scripts/02_build_manifest.py", font_size=14)
    add_bullet(tf2, "Logic: src/marine_peptides/download/biosample.py", font_size=14)

    add_bullet(tf2, "", font_size=8)
    add_bullet(tf2, "Technical notes:", font_size=14, bold=True, color=MID_GRAY)
    add_bullet(tf2, "SAMEA/SAMD need esearch→UID→efetch", font_size=13, color=MID_GRAY)
    add_bullet(tf2, "Results cached to data/interim/ (resumable)", font_size=13, color=MID_GRAY)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 6: Tier 1+2 Combined Results
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Tier 1 + 2 — Combined Results (GTDB R220 Marine)",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=28, color=NAVY)

    rows_data = [
        ["Metric", "Value"],
        ["Total genomes", "42,099"],
        ["GTDB species representatives", "21,062"],
        ["Bacteria", "38,328 (91%)"],
        ["Archaea", "3,771 (9%)"],
        ["MAGs", "34,420 (82%)"],
        ["Isolates", "6,383 (15%)"],
        ["SAGs", "1,296 (3%)"],
        ["Pass MIMAG QC (≥50% comp, ≤5% cont)", "40,061 (95.2%)"],
        ["With lat/lon coordinates", "20,989"],
        ["Host-associated (tagged)", "1,408"],
    ]
    add_table(slide, rows_data,
              [Inches(5.0), Inches(3.0)],
              Inches(0.7), Inches(1.4))

    tf = add_body_box(slide, Inches(0.7), Inches(6.0), Inches(12), Inches(1.0))
    add_bullet(tf, "Manifest: data/processed/genome_manifest.tsv (version-controlled provenance table)",
               font_size=14, color=MID_GRAY)

    # Top phyla on right side
    tf2 = add_body_box(slide, Inches(9.0), Inches(1.4), Inches(4.0), Inches(4.5))
    p = tf2.paragraphs[0]
    p.text = "Top Phyla (GTDB R220)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    phyla = [
        ("Pseudomonadota", "17,242"),
        ("Bacteroidota", "5,286"),
        ("Actinomycetota", "2,362"),
        ("Chloroflexota", "1,716"),
        ("Cyanobacteriota", "1,332"),
        ("Desulfobacterota", "1,158"),
        ("Planctomycetota", "1,147"),
        ("Thermoplasmatota", "1,136"),
    ]
    for name, count in phyla:
        add_bullet(tf2, f"{name}: {count}", font_size=12)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 7: Genome Download (Tier 1+2)
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Genome Download — Tier 1+2 (NCBI Datasets)",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=28, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "Method"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "All 42,099 genomes have NCBI GCA/GCF accessions (via GTDB)")
    add_bullet(tf, "Downloaded using NCBI Datasets CLI / REST API")
    add_bullet(tf, "FASTA-only (genomic sequences, no annotations)")
    add_bullet(tf, "Stored at: data/raw/ncbi_dataset/data/<accession>/")

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Storage", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "One directory per accession under data/raw/ncbi_dataset/data/")
    add_bullet(tf, "Each contains: <accession>_<assembly>_genomic.fna")
    add_bullet(tf, "Total: ~148 GB on disk")

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Key Design Choices", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "Manifest-driven: only download accessions in genome_manifest.tsv")
    add_bullet(tf, "Batched: 500 accessions per API call")
    add_bullet(tf, "data/raw/ is READ-ONLY after download (project convention)")

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 8: Tier 3 — Overview
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Tier 3 — External Curated Marine Catalogs",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=30, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.2), Inches(12), Inches(2.0))
    p = tf.paragraphs[0]
    p.text = "Why Tier 3?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "GTDB R220 only indexes NCBI GenBank/RefSeq — many marine genomes live elsewhere")
    add_bullet(tf, "Published after GTDB R220 freeze, or in non-NCBI repositories (ENA, Figshare, CNGB)")
    add_bullet(tf, "De novo assembled MAGs never submitted to GenBank individually")
    add_bullet(tf, "Goal: maximize recall by integrating all major marine genome catalogs")

    rows_data = [
        ["Catalog", "Source", "Type", "Genomes", "Disk", "Status"],
        ["MarRef", "Marine Metagenomics Portal v1.8", "References", "137", "151 MB", "Complete"],
        ["GORG-Tropics", "ENA PRJEB33281", "SAGs", "12,711", "2.4 GB", "Complete"],
        ["MarDB", "Marine Metagenomics Portal v1.7", "All types", "19,761", "38 GB", "Complete"],
        ["OceanDNA", "Figshare 5564844", "MAGs", "52,325", "31 GB", "Complete"],
        ["GOMC", "CNGB CNP0001755", "MAGs", "0", "0", "Failed (server down)"],
        ["TOTAL", "", "", "84,934", "~71 GB", ""],
    ]
    add_table(slide, rows_data,
              [Inches(1.5), Inches(3.5), Inches(1.2), Inches(1.2), Inches(1.2), Inches(2.0)],
              Inches(0.7), Inches(3.8))

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 9: Tier 3 — Download Strategy
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Tier 3 — Download Strategy & Design Principles",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=28, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "Design Principles"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "Resumable: state files track progress; safe to interrupt/restart")
    add_bullet(tf, "250 GB disk budget enforced (hard stop)")
    add_bullet(tf, "Smallest-first order: MarRef → GORG → MarDB → OceanDNA → GOMC")
    add_bullet(tf, "FASTA-only: no annotations or protein files")
    add_bullet(tf, "Symlink reuse: if genome already in Tier 1+2, symlink instead of re-download")
    add_bullet(tf, "Delete archives after extraction to save space")

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Parallel Downloads", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "NCBI REST batches: 4 workers × 100 accessions/batch")
    add_bullet(tf, "GORG SAGs: 12 parallel ENA FTP downloads")
    add_bullet(tf, "GOMC: 24-part HTTP Range parallel download (slow server)")
    add_bullet(tf, "File materialization: 8 parallel workers")

    tf2 = add_body_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf2.paragraphs[0]
    p.text = "Deduplication Strategy"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf2, "6,637 genomes exist in both Tier 1+2 and Tier 3")
    add_bullet(tf2, "These are stored as symlinks → Tier 1+2 files")
    add_bullet(tf2, "Zero disk duplication; manifest tracks is_in_tier12 flag")
    add_bullet(tf2, "")
    add_bullet(tf2, "File Organization", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf2, "data/raw/tier3/", font_size=14, bold=True)
    add_bullet(tf2, "  ├── marref/     137 files  (GCA_*.fna.gz)", font_size=13)
    add_bullet(tf2, "  ├── gorg/     12,711 files (<WGS>.fna.gz)", font_size=13)
    add_bullet(tf2, "  ├── mardb/    19,761 files (GCA_*_genomic.fna)", font_size=13)
    add_bullet(tf2, "  ├── oceandna/ 52,325 files (OceanDNA-*.fna.gz)", font_size=13)
    add_bullet(tf2, "  └── gomc/          0 files (failed)", font_size=13)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 10: Tier 3 — MarRef & MarDB
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Tier 3 — MarRef & MarDB (NCBI-backed catalogs)",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=28, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "Source: Marine Metagenomics Portal (mmp.sfb.uit.no)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "MarRef v1.8: 137 curated complete/near-complete marine reference genomes")
    add_bullet(tf, "MarDB v1.7: 19,761 comprehensive marine genomes (MAGs, SAGs, isolates)")
    add_bullet(tf, "All entries have NCBI GCA assembly accessions")

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Download Process", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "1. Parse metadata TSV → extract GCA accessions (acc:assembly_genbank column)")
    add_bullet(tf, "2. Check against Tier 1+2 local files → symlink matches (6,637 total)")
    add_bullet(tf, "3. Download remaining via NCBI Datasets REST API (batches of 100, 4 workers)")
    add_bullet(tf, "4. Extract zip → materialize .fna files to output dir (8 parallel workers)")
    add_bullet(tf, "5. Log unresolvable accessions to missing.txt")

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Metadata Enrichment (MarDB)", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "GTDB taxonomy available for 65% of genomes (tax:gtdb_classification)")
    add_bullet(tf, "CheckM completeness/contamination for 66% (gen:completeness, gen:contamination)")
    add_bullet(tf, "Isolation source, lat/lon, depth where recorded in source metadata")

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 11: Tier 3 — GORG-Tropics
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Tier 3 — GORG-Tropics (SAGs from ENA)",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=28, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(6.0), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "About"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "Global Ocean Reference Genomes — Tropics")
    add_bullet(tf, "Pachiadaki et al. 2019 (Nature)")
    add_bullet(tf, "Single-amplified genomes from tropical/subtropical surface ocean")
    add_bullet(tf, "ENA Project: PRJEB33281")
    add_bullet(tf, "12,711 SAG assemblies (WGS sets)")

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Download Process", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "1. Query ENA filereport API for PRJEB33281")
    add_bullet(tf, "   → get WGS set accessions + FTP URLs", level=1, font_size=14)
    add_bullet(tf, "2. Download each SAG .fasta.gz from ENA FTP")
    add_bullet(tf, "   12 parallel workers via ThreadPoolExecutor", level=1, font_size=14)
    add_bullet(tf, "3. Save as <WGS_SET_ID>.fna.gz")

    tf2 = add_body_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf2.paragraphs[0]
    p.text = "Technical Details"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf2, "No NCBI accession (ENA WGS sets, not in GenBank)")
    add_bullet(tf2, "No per-genome QC metadata available")
    add_bullet(tf2, "All tagged: isolation_source = 'marine water'")
    add_bullet(tf2, "")
    add_bullet(tf2, "ENA filereport URL format:", font_size=14, bold=True)
    add_bullet(tf2, "https://www.ebi.ac.uk/ena/portal/api/filereport", font_size=12)
    add_bullet(tf2, "  ?accession=PRJEB33281", font_size=12)
    add_bullet(tf2, "  &result=wgs_set&fields=...", font_size=12)
    add_bullet(tf2, "")
    add_bullet(tf2, "FTP pattern:", font_size=14, bold=True)
    add_bullet(tf2, "ftp.ebi.ac.uk/pub/databases/ena/wgs/", font_size=12)
    add_bullet(tf2, "  public/<prefix>/<SET_ID>.fasta.gz", font_size=12)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 12: Tier 3 — OceanDNA
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Tier 3 — OceanDNA (~52k Ocean MAGs from Figshare)",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=28, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(6.0), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "About"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "OceanDNA — Nishimura & Yoshizawa 2022")
    add_bullet(tf, "~52k MAGs from global ocean metagenomes")
    add_bullet(tf, "Hosted on Figshare (collection 5564844)")
    add_bullet(tf, "Includes species representatives + non-representatives")
    add_bullet(tf, "De novo assemblies — no NCBI accessions")

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Download Process", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "1. Download 3 tar archives from Figshare (direct file IDs):")
    add_bullet(tf, "   • fasta_species-representatives.tar (ID: 29164842)", level=1, font_size=14)
    add_bullet(tf, "   • fasta_non-representatives.tar (ID: 29200011)", level=1, font_size=14)
    add_bullet(tf, "   • Supplement_revised.tar.gz (ID: 35080240)", level=1, font_size=14)
    add_bullet(tf, "2. Verify MD5 checksums before extraction")
    add_bullet(tf, "3. Stream tar extraction → individual .fna.gz files")
    add_bullet(tf, "4. Delete source tar archives after extraction (save ~25 GB)")

    tf2 = add_body_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf2.paragraphs[0]
    p.text = "Technical Details"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf2, "52,325 genomes extracted")
    add_bullet(tf2, "31 GB total on disk")
    add_bullet(tf2, "No per-genome metadata (no GTDB, no CheckM)")
    add_bullet(tf2, "Naming: OceanDNA-<id>.fna.gz")
    add_bullet(tf2, "")
    add_bullet(tf2, "Figshare download URL pattern:", font_size=14, bold=True)
    add_bullet(tf2, "https://figshare.com/ndownloader/files/<FILE_ID>", font_size=12)
    add_bullet(tf2, "")
    add_bullet(tf2, "MD5 verification:", font_size=14, bold=True)
    add_bullet(tf2, "reps: 807b0f70c096ad740377e9e22c77b315", font_size=12)
    add_bullet(tf2, "nonreps: c50cd5c886441990ad46bf5cf001c6d0", font_size=12)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 13: Tier 3 — GOMC (failed)
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Tier 3 — GOMC (Failed — CNGB Server Unreachable)",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=28, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(6.0), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "About"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "Global Ocean Microbial Census (GOMC)")
    add_bullet(tf, "~67k MAGs from global ocean sampling")
    add_bullet(tf, "Hosted at CNGB (China National GeneBank)")
    add_bullet(tf, "Two archives: 24195 genomes (17.9 GB) + 43191 MAGs (30.8 GB)")

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "What Was Attempted", font_size=18, bold=True, color=ACCENT_ORANGE)
    add_bullet(tf, "1. Fetched md5.txt for integrity verification")
    add_bullet(tf, "2. Implemented 24-part parallel HTTP Range download")
    add_bullet(tf, "3. Multiple retry attempts over several days")
    add_bullet(tf, "4. Server consistently refused connections / SSL errors")

    tf2 = add_body_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf2.paragraphs[0]
    p.text = "Failure Mode"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED

    add_bullet(tf2, "Connection refused (Errno 111)")
    add_bullet(tf2, "SSL: UNEXPECTED_EOF_WHILE_READING")
    add_bullet(tf2, "Connection timeouts (134+ seconds)")
    add_bullet(tf2, "416 Range Not Satisfiable (corrupted partial)")
    add_bullet(tf2, "")
    add_bullet(tf2, "Future Retry", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf2, "Code is ready: scripts/04_download_tier3.py --catalog gomc")
    add_bullet(tf2, "Run when CNGB server is back online")
    add_bullet(tf2, "Script will automatically resume from scratch")
    add_bullet(tf2, "Expected: +67k genomes, +48 GB disk")

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 14: Combined Dataset
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Combined Dataset — All Tiers",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=30, color=NAVY)

    rows_data = [
        ["Tier", "Source", "Genomes", "Unique", "Disk"],
        ["1", "GTDB R220 keyword filter", "37,069", "37,069", "~100 GB"],
        ["2", "NCBI BioSample ENVO", "5,030", "5,030", "~48 GB"],
        ["3 - MarRef", "Marine Metagenomics Portal", "137", "121", "151 MB"],
        ["3 - GORG", "ENA (SAGs)", "12,711", "12,711", "2.4 GB"],
        ["3 - MarDB", "Marine Metagenomics Portal", "19,761", "13,140", "38 GB"],
        ["3 - OceanDNA", "Figshare (MAGs)", "52,325", "52,325", "31 GB"],
        ["3 - GOMC", "CNGB (failed)", "0", "0", "0"],
        ["TOTAL", "", "127,033", "~120,396", "~220 GB"],
    ]
    add_table(slide, rows_data,
              [Inches(1.5), Inches(3.5), Inches(1.5), Inches(1.5), Inches(1.5)],
              Inches(0.7), Inches(1.5))

    tf = add_body_box(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(1.5))
    add_bullet(tf, "6,637 genomes shared between Tier 1+2 and Tier 3 (stored as symlinks, no disk duplication)",
               font_size=14)
    add_bullet(tf, "78,297 genomes are unique to Tier 3 (not in GTDB marine filter results)", font_size=14)
    add_bullet(tf, "This is one of the most comprehensive marine prokaryote genome collections for peptide discovery",
               font_size=14, bold=True, color=OCEAN_BLUE)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 15: Code & Reproducibility
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "Code Organization & Reproducibility",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=28, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "Scripts (thin wrappers)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "00_select_marine_candidates.py → Tier 1", font_size=14)
    add_bullet(tf, "01_enrich_marine_biosample.py → Tier 2", font_size=14)
    add_bullet(tf, "02_build_manifest.py → combine Tier 1+2", font_size=14)
    add_bullet(tf, "03_fetch_tier3_metadata.py → Tier 3 metadata", font_size=14)
    add_bullet(tf, "04_download_tier3.py --catalog <name> → download", font_size=14)
    add_bullet(tf, "05_build_tier3_manifest.py → Tier 3 manifest", font_size=14)

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Package (reusable logic)", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "src/marine_peptides/download/", font_size=14, bold=True)
    add_bullet(tf, "  ├── gtdb_metadata.py", font_size=13)
    add_bullet(tf, "  ├── marine_filter.py", font_size=13)
    add_bullet(tf, "  ├── biosample.py", font_size=13)
    add_bullet(tf, "  ├── manifest.py", font_size=13)
    add_bullet(tf, "  └── tier3/", font_size=13)
    add_bullet(tf, "      ├── common.py (download, extract, disk)", font_size=12)
    add_bullet(tf, "      ├── marref.py (NCBI catalog logic)", font_size=12)
    add_bullet(tf, "      ├── gorg.py, oceandna.py, gomc.py", font_size=12)
    add_bullet(tf, "      └── manifest.py", font_size=12)

    tf2 = add_body_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf2.paragraphs[0]
    p.text = "Configuration"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf2, "config/config.yaml", font_size=14, bold=True)
    add_bullet(tf2, "  • All keyword lists (include/exclude/host)")
    add_bullet(tf2, "  • QC thresholds")
    add_bullet(tf2, "  • Tier 3 catalog configs (URLs, file IDs, MD5s)")
    add_bullet(tf2, "  • Disk budget (250 GB)")
    add_bullet(tf2, "  • Download order")
    add_bullet(tf2, "")
    add_bullet(tf2, "Manifests (outputs)", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf2, "data/processed/genome_manifest.tsv", font_size=14)
    add_bullet(tf2, "  → 42,099 rows (Tier 1+2 provenance)")
    add_bullet(tf2, "data/processed/tier3_manifest.tsv", font_size=14)
    add_bullet(tf2, "  → 84,934 rows (Tier 3 per-catalog inventory)")
    add_bullet(tf2, "")
    add_bullet(tf2, "To Reproduce", font_size=18, bold=True, color=OCEAN_BLUE)
    add_bullet(tf2, "1. Install: uv sync")
    add_bullet(tf2, "2. Run scripts 00→05 in order")
    add_bullet(tf2, "3. Everything regenerates from metadata + config")

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 16: How to Resume / Extend
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_box(slide, "For Future Contributors — How to Resume or Extend",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.8), font_size=28, color=NAVY)

    tf = add_body_box(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf.paragraphs[0]
    p.text = "Retry GOMC download"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf, "cd /path/to/marine", font_size=13)
    add_bullet(tf, "uv run python scripts/04_download_tier3.py --catalog gomc", font_size=13)
    add_bullet(tf, "# Then rebuild manifest:", font_size=13)
    add_bullet(tf, "uv run python scripts/05_build_tier3_manifest.py", font_size=13)

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Add a new catalog", font_size=16, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "1. Add config entry in config/config.yaml under tier3:", font_size=13)
    add_bullet(tf, "2. Create src/marine_peptides/download/tier3/<name>.py", font_size=13)
    add_bullet(tf, "3. Add handler in scripts/04_download_tier3.py", font_size=13)
    add_bullet(tf, "4. Re-run scripts/05_build_tier3_manifest.py", font_size=13)

    add_bullet(tf, "", font_size=8)
    add_bullet(tf, "Re-run full pipeline from scratch", font_size=16, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "# Ensure GTDB metadata in data/metadata/gtdb/", font_size=13)
    add_bullet(tf, "uv run python scripts/00_select_marine_candidates.py", font_size=13)
    add_bullet(tf, "uv run python scripts/01_enrich_marine_biosample.py", font_size=13)
    add_bullet(tf, "uv run python scripts/02_build_manifest.py", font_size=13)
    add_bullet(tf, "uv run python scripts/03_fetch_tier3_metadata.py", font_size=13)
    add_bullet(tf, "uv run python scripts/04_download_tier3.py", font_size=13)
    add_bullet(tf, "uv run python scripts/05_build_tier3_manifest.py", font_size=13)

    tf2 = add_body_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5))
    p = tf2.paragraphs[0]
    p.text = "Key Files to Know"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE

    add_bullet(tf2, "AGENTS.md — project conventions (read first!)", font_size=14)
    add_bullet(tf2, "config/config.yaml — all params in one place", font_size=14)
    add_bullet(tf2, "docs/marine-genome-selection.md — this documentation", font_size=14)
    add_bullet(tf2, "")
    add_bullet(tf2, "Important Constraints", font_size=16, bold=True, color=RED)
    add_bullet(tf2, "data/raw/ is READ-ONLY after download", font_size=14)
    add_bullet(tf2, "Never hardcode paths — use config.yaml", font_size=14)
    add_bullet(tf2, "Scripts are thin wrappers; logic in src/", font_size=14)
    add_bullet(tf2, "Disk budget: check before large downloads", font_size=14)
    add_bullet(tf2, "")
    add_bullet(tf2, "Environment Setup", font_size=16, bold=True, color=OCEAN_BLUE)
    add_bullet(tf2, "Python/ML: uv sync (reads pyproject.toml)", font_size=14)
    add_bullet(tf2, "Bioinfo CLIs: conda env create -f env/environment.yml", font_size=14)
    add_bullet(tf2, "Both needed for full pipeline", font_size=14)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 17: SmORFinder section divider
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_title_box(slide, "SmORFinder — Small ORF / Peptide Prediction",
                  Inches(1), Inches(2.2), Inches(11), Inches(1.0), font_size=34, color=WHITE)
    add_title_box(slide, "Hybrid Nextflow run on the full genome collection",
                  Inches(1), Inches(3.3), Inches(11), Inches(0.7), font_size=20, bold=False, color=LIGHT_BLUE)
    tf = add_body_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1.5))
    add_bullet(tf, "112,414 genomes  →  112,097 completed (99.72%)  →  474,562 predicted smORFs",
               font_size=18, color=WHITE)
    add_bullet(tf, "Finished 2026-07-14  ·  docs/smorfinder-pipeline.md", font_size=14, color=MID_GRAY)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 18: SmORFinder — how executed
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_box(slide, "SmORFinder — How the Run Was Executed",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.7), font_size=28, color=NAVY)

    flow = [
        ("1. Build inputs\ngenome_index.tsv\n→ smorfinder_inputs.tsv", Inches(0.5), NAVY),
        ("2. Phase 1: single\nNextflow + SLURM\nsmorf single × 112k", Inches(3.6), OCEAN_BLUE),
        ("3. Promote failures\n593 genomes → meta\n(script 13)", Inches(6.7), TEAL),
        ("4. Phase 2: meta\nre-run only failures\nsmorf meta", Inches(9.8), DARK_TEAL),
    ]
    for text, left, color in flow:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.2), Inches(2.9), Inches(1.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13) if i > 0 else Pt(15)
            p.font.bold = i == 0
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    tf = add_body_box(slide, Inches(0.7), Inches(2.8), Inches(6.0), Inches(4.2))
    p = tf.paragraphs[0]
    p.text = "Execution details"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE
    add_bullet(tf, "Workflow: smorfinder.nf  (profile: slurm)", font_size=14)
    add_bullet(tf, "Env: dedicated conda (Python 3.8 / TF 2.3.1)", font_size=14)
    add_bullet(tf, "Concurrency: 12 parallel tasks", font_size=14)
    add_bullet(tf, "Resumable: Nextflow -resume + per-genome _SUCCESS.json", font_size=14)
    add_bullet(tf, "Failed genomes ignored (not fatal) → promoted to meta", font_size=14)
    add_bullet(tf, "Why hybrid: single-mode Prodigal segfaults on ~0.5%", font_size=14)
    add_bullet(tf, "  of genomes; meta recovers most of them", level=1, font_size=13)

    tf2 = add_body_box(slide, Inches(7.0), Inches(2.8), Inches(5.8), Inches(4.2))
    p = tf2.paragraphs[0]
    p.text = "Key commands"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE
    add_bullet(tf2, "scripts/12_build_smorfinder_inputs_from_index.py", font_size=13)
    add_bullet(tf2, "nextflow run smorfinder.nf -profile slurm -resume", font_size=13)
    add_bullet(tf2, "scripts/13_promote_smorfinder_failures.py", font_size=13)
    add_bullet(tf2, "nextflow run smorfinder.nf -profile slurm -resume", font_size=13)
    add_bullet(tf2, "", font_size=6)
    add_bullet(tf2, "Wall time", font_size=16, bold=True, color=OCEAN_BLUE)
    add_bullet(tf2, "Phase 1 (single): ~1d 9h", font_size=14)
    add_bullet(tf2, "Phase 2 (meta): ~12 min", font_size=14)
    add_bullet(tf2, "Log: logs/smorfinder-hybrid-20260713-133555.log", font_size=12, color=MID_GRAY)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 19: SmORFinder — outputs
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_box(slide, "SmORFinder — Results & Output Locations",
                  Inches(0.7), Inches(0.3), Inches(12), Inches(0.7), font_size=28, color=NAVY)

    rows_data = [
        ["Metric", "Count"],
        ["Input genomes", "112,414"],
        ["single-mode successes", "111,821"],
        ["meta-mode recoveries", "276"],
        ["Hard failures (mostly GORG SAGs)", "317"],
        ["Coverage", "112,097 (99.72%)"],
        ["Predicted smORFs", "474,562"],
    ]
    add_table(slide, rows_data, [Inches(4.2), Inches(2.2)], Inches(0.7), Inches(1.15))

    tf = add_body_box(slide, Inches(7.3), Inches(1.15), Inches(5.5), Inches(3.6))
    p = tf.paragraphs[0]
    p.text = "Per-genome outputs"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE
    add_bullet(tf, "data/interim/smorfinder/single/<id>/", font_size=13, bold=True)
    add_bullet(tf, "data/interim/smorfinder/meta/<id>/", font_size=13, bold=True)
    add_bullet(tf, "Each dir: .faa  .ffn  .gff  .tsv  _SUCCESS.json", font_size=13)
    add_bullet(tf, "", font_size=6)
    add_bullet(tf, "Run summary (authoritative)", font_size=16, bold=True, color=OCEAN_BLUE)
    add_bullet(tf, "data/processed/smorfinder_run_manifest.tsv", font_size=13, bold=True)
    add_bullet(tf, "", font_size=6)
    add_bullet(tf, "Hard failures list", font_size=16, bold=True, color=ACCENT_ORANGE)
    add_bullet(tf, "results/tables/smorfinder_hard_failures.tsv", font_size=13)

    tf2 = add_body_box(slide, Inches(0.7), Inches(5.1), Inches(12), Inches(2.0))
    p = tf2.paragraphs[0]
    p.text = "What each file is"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = OCEAN_BLUE
    add_bullet(tf2, ".faa / .ffn — predicted small proteins (AA) and coding sequences (nt)", font_size=14)
    add_bullet(tf2, ".gff / .tsv — coordinates + scores (HMM + deep models); _SUCCESS.json — resume marker", font_size=14)
    add_bullet(tf2, "Docs: docs/smorfinder-pipeline.md  ·  No re-run needed unless inputs/cutoffs change",
               font_size=14, color=GREEN)

    # ═══════════════════════════════════════════════════════════════════════════
    # Save
    # ═══════════════════════════════════════════════════════════════════════════
    out_path = "docs/marine_genome_download_pipeline.pptx"
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    create_presentation()
